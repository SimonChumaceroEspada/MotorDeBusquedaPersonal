import os
import sys
import io
import PyPDF2
from docx import Document as DocxDocument
import traceback

# Importaciones para Excel y PowerPoint
try:
    import openpyxl
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False
    print("Warning: openpyxl not installed. Excel support will be limited.")

try:
    from pptx import Presentation
    PPT_SUPPORT = True
except ImportError:
    PPT_SUPPORT = False
    print("Warning: python-pptx not installed. PowerPoint support will be limited.")

def process_documents(docs_dir):
    """Process all documents in the given directory and return their content"""
    documents = []
    
    # Check if directory exists
    if not os.path.exists(docs_dir):
        print(f"Documents directory not found: {docs_dir}")
        return documents
    
    # Walk through all files in the directory
    for root, _, files in os.walk(docs_dir):
        for filename in files:
            filepath = os.path.join(root, filename)
            extension = os.path.splitext(filename)[1].lower()
            
            try:
                # Process file based on extension
                content = ""
                
                if extension == '.txt':
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                
                elif extension == '.pdf':
                    content = extract_pdf_content(filepath)
                
                elif extension in ['.docx', '.doc']:
                    content = extract_docx_content(filepath)
                
                elif extension in ['.xlsx', '.xls']:
                    content = extract_excel_content(filepath)
                
                elif extension in ['.pptx', '.ppt']:
                    content = extract_ppt_content(filepath)
                
                # Skip empty or unprocessable files
                if not content or len(content.strip()) == 0:
                    print(f"Skipping {filename}: No content extracted")
                    continue
                
                # Add to documents list
                documents.append({
                    'filename': filename,
                    'path': filepath,
                    'extension': extension,
                    'content': content
                })
                
                print(f"Processed {filename}: {len(content)} characters")
                
            except Exception as e:
                print(f"Error processing {filename}: {str(e)}")
                traceback.print_exc()
    
    return documents

def extract_pdf_content(filepath):
    """Extract text content from PDF file"""
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            content = ""
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                content += page.extract_text() + "\n"
            return content
    except Exception as e:
        print(f"PDF extraction error: {str(e)}")
        return ""

def extract_docx_content(filepath):
    """Extract text content from DOCX file"""
    try:
        doc = DocxDocument(filepath)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"DOCX extraction error: {str(e)}")
        return ""

def extract_excel_content(filepath):
    """Extract text content from Excel file"""
    if not EXCEL_SUPPORT:
        print("Excel support not available. Please install openpyxl.")
        return ""
    
    try:
        workbook = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        full_text = []
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            # Iterate through all cells with data
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value is not None:
                        row_text.append(str(cell.value))
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            full_text.append("\n".join(sheet_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        print(f"Excel extraction error: {str(e)}")
        traceback.print_exc()
        return ""

def extract_ppt_content(filepath):
    """Extract text content from PowerPoint file"""
    if not PPT_SUPPORT:
        print("PowerPoint support not available. Please install python-pptx.")
        return ""
    
    try:
        presentation = Presentation(filepath)
        full_text = []
        
        # Process title and subtitle on each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = [f"Slide {i+1}:"]
            
            # Get text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                
                # Check for tables
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if len(slide_text) > 1:  # If there's more than just the slide number
                full_text.append("\n".join(slide_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        print(f"PowerPoint extraction error: {str(e)}")
        traceback.print_exc()
        return ""

if __name__ == "__main__":
    # Test file processing
    if len(sys.argv) > 1:
        test_dir = sys.argv[1]
    else:
        test_dir = "../documents"
    
    results = process_documents(test_dir)
    print(f"Processed {len(results)} documents")
