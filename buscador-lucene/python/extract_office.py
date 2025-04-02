#!/usr/bin/env python
"""
Script para extraer contenido de archivos Excel y PowerPoint
Utilizado por el backend Node.js para procesar estos formatos
No depende de otras dependencias como python-docx
"""

import sys
import os
import traceback
import json

# Importaciones individuales con manejo de errores
try:
    import openpyxl
except ImportError:
    print("Error: Falta instalar la dependencia 'openpyxl'")
    print("Por favor, ejecuta: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("Error: Falta instalar la dependencia 'python-pptx'")
    print("Por favor, ejecuta: pip install python-pptx")
    sys.exit(1)

def extract_excel_content(filepath):
    """Extract text content from Excel file"""
    try:
        workbook = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        full_text = []
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            # Iterate through all cells with data
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value is not None:
                        row_text.append(str(cell.value))
                if row_text:
                    sheet_text.append(" | ".join(row_text))
            
            full_text.append("\n".join(sheet_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        print(f"Excel extraction error: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        return ""

def extract_ppt_content(filepath):
    """Extract text content from PowerPoint file"""
    try:
        presentation = Presentation(filepath)
        full_text = []
        
        # Process title and subtitle on each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = [f"Slide {i+1}:"]
            
            # Get text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                
                # Check for tables
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if len(slide_text) > 1:  # If there's more than just the slide number
                full_text.append("\n".join(slide_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        print(f"PowerPoint extraction error: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        return ""

def find_query_in_content(content, query, context_size=50):
    """
    Busca una consulta en el contenido y devuelve fragmentos con contexto
    
    Args:
        content: El texto completo donde buscar
        query: La consulta a buscar
        context_size: Número de caracteres de contexto antes/después de la coincidencia
    
    Returns:
        Dict con el contenido completo y una lista de fragmentos con contexto
    """
    results = {
        "content": content,
        "matches": []
    }
    
    content_lower = content.lower()
    query_lower = query.lower()
    start_pos = 0
    found_pos = content_lower.find(query_lower, start_pos)
    
    while found_pos != -1:
        # Calcular el inicio y fin del fragmento con contexto
        start_extract = max(0, found_pos - context_size)
        end_extract = min(len(content), found_pos + len(query_lower) + context_size)
        
        # Extraer fragmento con contexto
        excerpt = content[start_extract:end_extract]
        
        # Guardar la coincidencia con su fragmento
        results["matches"].append({
            "position": found_pos,
            "excerpt": excerpt,
            "location": get_location_info(content, found_pos)
        })
        
        # Buscar la siguiente ocurrencia
        start_pos = found_pos + len(query_lower)
        found_pos = content_lower.find(query_lower, start_pos)
    
    return results

def get_location_info(content, position):
    """
    Identifica en qué hoja/diapositiva se encuentra la coincidencia
    
    Args:
        content: El contenido completo
        position: La posición de la coincidencia
    
    Returns:
        Información de ubicación (hoja/diapositiva)
    """
    # Encontrar la última marca de hoja/diapositiva antes de la posición
    lines = content[:position].split('\n')
    
    for i in range(len(lines)-1, -1, -1):
        if lines[i].startswith('Sheet:'):
            return lines[i]
        elif lines[i].startswith('Slide '):
            return lines[i]
    
    return "Ubicación desconocida"

def main():
    """
    Extrae el contenido de un archivo Excel o PowerPoint
    Uso: python extract_office.py <ruta_archivo> <extensión> [<consulta>]
    """
    if len(sys.argv) < 3:
        print("Error: Se requieren al menos 2 argumentos")
        print("Uso: python extract_office.py <ruta_archivo> <extensión> [<consulta>]")
        sys.exit(1)
    
    file_path = sys.argv[1]
    extension = sys.argv[2].lower()
    query = sys.argv[3].lower() if len(sys.argv) > 3 else None
    
    if not os.path.exists(file_path):
        print(f"Error: No se encontró el archivo {file_path}")
        sys.exit(1)
    
    try:
        content = ""
        if extension in ['.xlsx', '.xls']:
            content = extract_excel_content(file_path)
        elif extension in ['.pptx', '.ppt']:
            content = extract_ppt_content(file_path)
        else:
            print(f"Error: Formato no soportado - {extension}")
            sys.exit(1)
        
        # Si hay una consulta, buscar ocurrencias y devolver contextos
        if query and query.strip():
            results = find_query_in_content(content, query)
            print(json.dumps(results))
        else:
            # De lo contrario, simplemente devolver todo el contenido
            print(content)
            
        sys.exit(0)
    except Exception as e:
        print(f"Error al procesar archivo: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
